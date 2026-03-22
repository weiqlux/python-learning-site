#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档分析 + 思维导图生成器
支持：财务报表、流程图、架构图、思维导图等文档的智能分析和可视化
"""

import os
import sys
import json
import base64
import hashlib
import shutil
import http.client
import re
import logging
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import sqlite3
from PIL import Image
import io

# 尝试导入 xmind 库，如果失败则使用备用方案
try:
    import xmind
    XMind_AVAILABLE = True
except ImportError:
    XMind_AVAILABLE = False
    logger = logging.getLogger(__name__)
    logger.warning("xmind 库未安装，将使用备用方案生成 .xmind 文件")

# 尝试导入 python-pptx 库
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx 库未安装，PPT 导出功能将不可用")

# 配置日志
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout),
        logging.FileHandler('/tmp/document_mindmap.log', encoding='utf-8')
    ]
)
logger = logging.getLogger(__name__)


class DocumentMindMapGenerator:
    """文档分析与思维导图生成器"""
    
    # 预设的分析场景
    SCENARIOS = {
        "financial_report": {
            "name": "财务报表分析",
            "description": "分析财务报表中的表格数据、增长趋势、关键指标",
            "system_prompt": """你是一位专业的财务分析师。请分析用户上传的财务报表图片。

请完成以下任务：
1. 提取表格中的所有数据
2. 计算关键财务指标（同比增长率、环比增长率等）
3. 识别异常数据和趋势
4. 生成结构化的分析结果

输出格式必须是严格的 JSON：
{
    "document_type": "financial_report",
    "title": "报表标题",
    "tables": [
        {
            "table_name": "表格名称",
            "headers": ["列1", "列2", "列3"],
            "rows": [["数据1", "数据2", "数据3"], ...]
        }
    ],
    "key_metrics": [
        {"name": "指标名称", "value": "数值", "change": "变化率", "trend": "up/down/stable"}
    ],
    "insights": ["洞察1", "洞察2"],
    "mindmap_structure": {
        "root": "报表主题",
        "children": [
            {
                "name": "一级节点",
                "children": [
                    {"name": "二级节点"},
                    {"name": "二级节点"}
                ]
            }
        ]
    }
}"""
        },
        "flowchart": {
            "name": "流程图解读",
            "description": "解读流程图逻辑、分支条件、执行路径",
            "system_prompt": """你是一位系统架构师。请分析用户上传的流程图。

请完成以下任务：
1. 识别流程图的所有节点和步骤
2. 分析流程的分支逻辑和判断条件
3. 识别关键路径和异常处理
4. 解释特定场景下的执行流程

输出格式必须是严格的 JSON：
{
    "document_type": "flowchart",
    "title": "流程图标题",
    "nodes": [
        {"id": "1", "type": "start/action/decision/end", "label": "节点描述"}
    ],
    "edges": [
        {"from": "1", "to": "2", "label": "条件/动作"}
    ],
    "paths": [
        {"name": "主流程", "steps": ["步骤1", "步骤2"]}
    ],
    "key_decisions": [
        {"condition": "判断条件", "true_path": "是分支", "false_path": "否分支"}
    ],
    "mindmap_structure": {
        "root": "流程主题",
        "children": [
            {
                "name": "开始/输入",
                "children": [{"name": "输入项1"}, {"name": "输入项2"}]
            },
            {
                "name": "处理逻辑",
                "children": [{"name": "判断条件"}, {"name": "处理步骤"}]
            },
            {
                "name": "输出/结果",
                "children": [{"name": "成功"}, {"name": "失败/异常"}]
            }
        ]
    }
}"""
        },
        "architecture": {
            "name": "架构图分析",
            "description": "分析系统架构图、组件关系、数据流向",
            "system_prompt": """你是一位资深系统架构师。请分析用户上传的系统架构图。

请完成以下任务：
1. 识别所有系统组件和模块
2. 分析组件之间的关系和依赖
3. 识别数据流向和接口
4. 评估架构特点

输出格式必须是严格的 JSON：
{
    "document_type": "architecture",
    "title": "架构图标题",
    "components": [
        {"name": "组件名", "type": "service/db/cache/gateway", "description": "功能描述"}
    ],
    "relationships": [
        {"from": "组件A", "to": "组件B", "type": "sync/async", "protocol": "HTTP/GRPC/..."}
    ],
    "data_flow": ["数据源", "处理", "存储", "展示"],
    "architecture_patterns": ["微服务", "分层架构", "事件驱动"],
    "mindmap_structure": {
        "root": "系统架构",
        "children": [
            {
                "name": "接入层",
                "children": [{"name": "API Gateway"}, {"name": "负载均衡"}]
            },
            {
                "name": "服务层",
                "children": [{"name": "业务服务A"}, {"name": "业务服务B"}]
            },
            {
                "name": "数据层",
                "children": [{"name": "数据库"}, {"name": "缓存"}]
            }
        ]
    }
}"""
        },
        "mindmap": {
            "name": "思维导图归纳",
            "description": "归纳思维导图要点、层级结构、核心主题",
            "system_prompt": """你是一位知识管理专家。请分析用户上传的思维导图。

请完成以下任务：
1. 提取思维导图的核心主题
2. 识别所有层级和分支结构
3. 归纳每个分支的要点
4. 发现知识点之间的关联

输出格式必须是严格的 JSON：
{
    "document_type": "mindmap",
    "title": "导图主题",
    "central_theme": "中心主题",
    "main_branches": [
        {
            "name": "主分支1",
            "key_points": ["要点1", "要点2"],
            "sub_branches": [
                {"name": "子分支1", "points": ["子要点"]},
                {"name": "子分支2", "points": ["子要点"]}
            ]
        }
    ],
    "summary": "整体总结",
    "mindmap_structure": {
        "root": "中心主题",
        "children": [
            {
                "name": "主分支1",
                "children": [
                    {"name": "子分支1"},
                    {"name": "子分支2"}
                ]
            }
        ]
    }
}"""
        },
        "general": {
            "name": "通用文档分析",
            "description": "分析任意文档，提取关键信息",
            "system_prompt": """你是一位文档分析专家。请分析用户上传的文档图片。

请完成以下任务：
1. 识别文档类型和内容结构
2. 提取关键信息和数据
3. 总结核心观点
4. 生成思维导图结构

输出格式必须是严格的 JSON：
{
    "document_type": "general",
    "title": "文档标题",
    "document_category": "文档类别",
    "key_content": ["要点1", "要点2", "要点3"],
    "data_extracted": [
        {"type": "table/text/chart", "content": "具体内容"}
    ],
    "summary": "文档总结",
    "mindmap_structure": {
        "root": "文档主题",
        "children": [
            {"name": "要点1", "children": [{"name": "细节1"}]},
            {"name": "要点2", "children": [{"name": "细节2"}]}
        ]
    }
}"""
        }
    }
    
    def __init__(self, config: Dict = None):
        """初始化生成器"""
        self.config = config or {}
        self.api_key = self.config.get('api_key') or os.environ.get('DASHSCOPE_API_KEY', '')
        self.model = self.config.get('model', 'qwen-vl-max')
        
        # 目录设置
        self.base_dir = Path(__file__).parent
        self.upload_dir = self.base_dir / 'static' / 'uploads' / 'documents'
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # 初始化数据库
        self._init_db()
    
    def _init_db(self):
        """初始化数据库表"""
        conn = sqlite3.connect(str(self.base_dir / 'document_analysis.db'))
        cursor = conn.cursor()
        
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS document_sessions (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT UNIQUE NOT NULL,
                document_type TEXT,
                document_title TEXT,
                image_path TEXT,
                analysis_result TEXT,
                mindmap_data TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT NOT NULL,
                role TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (session_id) REFERENCES document_sessions(session_id)
            )
        ''')
        
        conn.commit()
        conn.close()
    
    def _get_db(self):
        """获取数据库连接"""
        conn = sqlite3.connect(str(self.base_dir / 'document_analysis.db'))
        conn.row_factory = sqlite3.Row
        return conn
    
    def compress_image(self, image_path: str, max_size: int = 800, quality: int = 75) -> bytes:
        """
        压缩图片以减少传输大小
        
        Args:
            image_path: 图片路径
            max_size: 最大边长（像素）
            quality: JPEG 质量（1-95）
            
        Returns:
            压缩后的图片 bytes
        """
        logger.info(f"开始压缩图片: {image_path}")
        
        # 获取原图信息
        original_size = os.path.getsize(image_path)
        logger.info(f"原图大小: {original_size / 1024:.2f} KB")
        
        with Image.open(image_path) as img:
            logger.info(f"图片尺寸: {img.size}, 模式: {img.mode}")
            # 转换为 RGB（处理 PNG 透明通道）
            if img.mode in ('RGBA', 'LA', 'P'):
                logger.info(f"转换图片模式: {img.mode} -> RGB")
                img = img.convert('RGB')
            
            # 等比例缩放
            original_dimensions = img.size
            # 兼容不同 PIL 版本
            try:
                resampling = Image.Resampling.LANCZOS
            except AttributeError:
                resampling = Image.ANTIALIAS  # 旧版本使用 ANTIALIAS
            img.thumbnail((max_size, max_size), resampling)
            logger.info(f"缩放图片: {original_dimensions} -> {img.size}")
            
            # 保存到内存
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=quality, optimize=True)
            compressed_bytes = buffer.getvalue()
            compressed_size = len(compressed_bytes)
            logger.info(f"压缩后大小: {compressed_size / 1024:.2f} KB (压缩率: {compressed_size/original_size*100:.1f}%)")
            
            return compressed_bytes
    
    def save_image(self, image_path: str) -> str:
        """保存图片到上传目录"""
        src = Path(image_path)
        if not src.exists():
            raise FileNotFoundError(f"图片文件不存在：{image_path}")
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        file_hash = hashlib.md5(src.read_bytes()).hexdigest()[:8]
        ext = src.suffix.lower()
        filename = f"doc-{timestamp}-{file_hash}{ext}"
        
        date_dir = datetime.now().strftime('%Y-%m')
        dest_dir = self.upload_dir / date_dir
        dest_dir.mkdir(parents=True, exist_ok=True)
        
        dest_path = dest_dir / filename
        shutil.copy2(src, dest_path)
        
        return f"static/uploads/documents/{date_dir}/{filename}"
    
    def analyze_document(
        self, 
        image_path: str, 
        question: str = None,
        scenario: str = "general",
        session_id: str = None
    ) -> Dict:
        """
        分析文档图片
        
        Args:
            image_path: 图片路径
            question: 用户问题（可选）
            scenario: 分析场景
            session_id: 会话ID（多轮对话用）
            
        Returns:
            分析结果
        """
        logger.info(f"开始分析文档: image_path={image_path}, scenario={scenario}, question={question}")
        
        # 压缩并编码图片
        try:
            compressed_image = self.compress_image(image_path)
            image_base64 = base64.b64encode(compressed_image).decode('utf-8')
            logger.info(f"图片 base64 编码完成, 长度: {len(image_base64)} 字符")
        except Exception as e:
            logger.error(f"图片压缩失败: {e}, 尝试使用原图")
            # 如果压缩失败，使用原图
            with open(image_path, 'rb') as f:
                image_base64 = base64.b64encode(f.read()).decode('utf-8')
            logger.info(f"原图 base64 编码完成, 长度: {len(image_base64)} 字符")
        
        # 获取场景配置
        scenario_config = self.SCENARIOS.get(scenario, self.SCENARIOS["general"])
        system_prompt = scenario_config["system_prompt"]
        
        # 构建用户问题
        if question:
            user_prompt = f"""
{system_prompt}

用户问题：{question}

请根据图片内容和用户问题进行分析，返回 JSON 格式结果。"""
        else:
            user_prompt = system_prompt
        
        # 构建请求
        request_body = json.dumps({
            "model": self.model,
            "input": {
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "image": f"data:image/jpeg;base64,{image_base64}"
                            },
                            {
                                "text": user_prompt
                            }
                        ]
                    }
                ]
            },
            "parameters": {
                "result_format": "message"
            }
        })
        
        request_body_size = len(request_body.encode('utf-8'))
        logger.info(f"请求体大小: {request_body_size / 1024:.2f} KB")
        logger.debug(f"请求体内容: {request_body[:500]}...")
        
        # 发送请求
        logger.info("开始发送 API 请求到 dashscope.aliyuncs.com")
        logger.info(f"API Key 状态: {'已配置' if self.api_key else '未配置'}, 长度: {len(self.api_key) if self.api_key else 0}")
        logger.debug(f"API Key: {self.api_key[:20]}..." if self.api_key else "API Key: None")
        
        try:
            conn = http.client.HTTPSConnection("dashscope.aliyuncs.com", timeout=60)
            headers = {
                'Authorization': f'Bearer {self.api_key}',
                'Content-Type': 'application/json'
            }
            
            logger.info(f"请求头: {headers}")
            logger.info("建立 HTTPS 连接...")
            conn.request(
                "POST",
                "/api/v1/services/aigc/multimodal-generation/generation",
                body=request_body,
                headers=headers
            )
            
            logger.info("等待 API 响应...")
            response = conn.getresponse()
            logger.info(f"收到响应: status={response.status}, reason={response.reason}")
            
            result = response.read().decode('utf-8')
            logger.info(f"响应内容长度: {len(result)} 字符")
            logger.debug(f"响应内容: {result[:500]}...")
            
            conn.close()
            logger.info("连接已关闭")
            
            # 解析响应
            logger.info("解析 API 响应...")
            result_json = json.loads(result)
            
            if 'output' in result_json and 'choices' in result_json['output']:
                content = result_json['output']['choices'][0]['message']['content']
                logger.info(f"API 返回内容类型: {type(content)}")
                
                if isinstance(content, list):
                    text = content[0].get('text', '') if content else ''
                else:
                    text = content
                
                logger.info(f"提取的文本长度: {len(text)} 字符")
                logger.debug(f"提取的文本内容: {text[:500]}...")
                
                # 提取 JSON
                analysis = self._extract_json(text)
                
                if 'error' in analysis:
                    logger.error(f"JSON 解析失败: {analysis['error']}")
                    return {
                        'success': False,
                        'error': f'JSON 解析失败: {analysis["error"]}',
                        'raw_text': analysis.get('raw_text', '')
                    }
                
                logger.info(f"JSON 解析成功, 包含字段: {list(analysis.keys())}")
                
                # 保存会话
                if not session_id:
                    session_id = f"session_{datetime.now().strftime('%Y%m%d%H%M%S')}_{hashlib.md5(image_path.encode()).hexdigest()[:6]}"
                
                saved_path = self.save_image(image_path)
                self._save_session(session_id, scenario, saved_path, analysis, question)
                logger.info(f"会话已保存: session_id={session_id}")
                
                return {
                    'success': True,
                    'session_id': session_id,
                    'analysis': analysis,
                    'mindmap_data': analysis.get('mindmap_structure', {})
                }
            else:
                logger.error(f"API 响应格式异常: {result_json.keys()}")
                return {
                    'success': False,
                    'error': 'API 响应格式异常',
                    'raw_response': result_json
                }
                
        except http.client.HTTPException as e:
            logger.error(f"HTTP 请求异常: {type(e).__name__}: {e}")
            return {
                'success': False,
                'error': f'HTTP 请求失败: {str(e)}'
            }
        except json.JSONDecodeError as e:
            logger.error(f"JSON 解析异常: {e}")
            return {
                'success': False,
                'error': f'响应解析失败: {str(e)}'
            }
        except Exception as e:
            logger.exception(f"未知异常: {type(e).__name__}: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def _extract_json(self, text: str) -> Dict:
        """从文本中提取 JSON"""
        # 尝试匹配 markdown 代码块
        json_match = re.search(r'```json\s*(.+?)\s*```', text, re.DOTALL)
        if json_match:
            text = json_match.group(1)
        else:
            # 尝试匹配普通 JSON
            json_match = re.search(r'\{[\s\S]*\}', text)
            if json_match:
                text = json_match.group(0)
        
        try:
            return json.loads(text)
        except json.JSONDecodeError as e:
            # 尝试修复常见 JSON 错误
            try:
                # 去除尾部逗号
                text = re.sub(r',(\s*[}\]])', r'\1', text)
                return json.loads(text)
            except:
                return {
                    'error': f'JSON 解析失败: {e}',
                    'raw_text': text[:2000]
                }
    
    def _save_session(
        self, 
        session_id: str, 
        document_type: str, 
        image_path: str, 
        analysis: Dict,
        question: str = None
    ):
        """保存会话到数据库"""
        conn = self._get_db()
        cursor = conn.cursor()
        
        # 检查是否已存在
        cursor.execute('SELECT id FROM document_sessions WHERE session_id = ?', (session_id,))
        existing = cursor.fetchone()
        
        if existing:
            cursor.execute('''
                UPDATE document_sessions 
                SET analysis_result = ?, mindmap_data = ?, updated_at = CURRENT_TIMESTAMP
                WHERE session_id = ?
            ''', (json.dumps(analysis, ensure_ascii=False), 
                  json.dumps(analysis.get('mindmap_structure', {}), ensure_ascii=False),
                  session_id))
        else:
            cursor.execute('''
                INSERT INTO document_sessions 
                (session_id, document_type, document_title, image_path, analysis_result, mindmap_data)
                VALUES (?, ?, ?, ?, ?, ?)
            ''', (session_id, document_type, 
                  analysis.get('title', '未命名文档'),
                  image_path,
                  json.dumps(analysis, ensure_ascii=False),
                  json.dumps(analysis.get('mindmap_structure', {}), ensure_ascii=False)))
        
        # 保存对话历史
        if question:
            cursor.execute('''
                INSERT INTO chat_history (session_id, role, content)
                VALUES (?, ?, ?)
            ''', (session_id, 'user', question))
            
            cursor.execute('''
                INSERT INTO chat_history (session_id, role, content)
                VALUES (?, ?, ?)
            ''', (session_id, 'assistant', json.dumps(analysis, ensure_ascii=False)))
        
        conn.commit()
        conn.close()
    
    def chat(
        self, 
        session_id: str, 
        question: str,
        image_path: str = None
    ) -> Dict:
        """
        多轮对话
        
        Args:
            session_id: 会话ID
            question: 用户问题
            image_path: 图片路径（可选，新图片）
            
        Returns:
            回答结果
        """
        # 获取历史会话
        conn = self._get_db()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT * FROM document_sessions WHERE session_id = ?
        ''', (session_id,))
        session = cursor.fetchone()
        
        if not session:
            conn.close()
            return {'success': False, 'error': '会话不存在'}
        
        # 获取历史对话
        cursor.execute('''
            SELECT role, content FROM chat_history 
            WHERE session_id = ? ORDER BY created_at
        ''', (session_id,))
        history = cursor.fetchall()
        conn.close()
        
        # 构建上下文
        context = []
        for h in history[-5:]:  # 最近5轮
            role, content = h
            if role == 'assistant':
                try:
                    content_json = json.loads(content)
                    content = content_json.get('summary', str(content_json)[:500])
                except:
                    pass
            context.append(f"{role}: {content}")
        
        context_str = "\n".join(context)
        
        # 获取图片并压缩
        img_path = image_path or session['image_path']
        if img_path and not img_path.startswith('/'):
            img_path = str(self.base_dir / img_path)
        
        try:
            compressed_image = self.compress_image(img_path)
            image_base64 = base64.b64encode(compressed_image).decode('utf-8')
        except Exception as e:
            # 如果压缩失败，使用原图
            with open(img_path, 'rb') as f:
                image_base64 = base64.b64encode(f.read()).decode('utf-8')
        
        # 构建提示词
        prompt = f"""基于之前的文档分析，回答用户的新问题。

历史对话：
{context_str}

用户新问题：{question}

请直接回答用户的问题，可以引用文档中的具体数据或信息。如果问题涉及新的分析，请返回 JSON 格式，包含：
{{
    "answer": "直接回答",
    "referenced_data": ["引用的数据1", "数据2"],
    "additional_insights": ["额外洞察"]
}}"""
        
        # 发送请求
        request_body = json.dumps({
            "model": self.model,
            "input": {
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "image": f"data:image/jpeg;base64,{image_base64}"
                            },
                            {
                                "text": prompt
                            }
                        ]
                    }
                ]
            },
            "parameters": {
                "result_format": "message"
            }
        })
        
        try:
            conn_http = http.client.HTTPSConnection("dashscope.aliyuncs.com")
            headers = {
                'Authorization': f'Bearer {self.api_key}',
                'Content-Type': 'application/json'
            }
            
            conn_http.request(
                "POST",
                "/api/v1/services/aigc/multimodal-generation/generation",
                body=request_body,
                headers=headers
            )
            
            response = conn_http.getresponse()
            result = response.read().decode('utf-8')
            conn_http.close()
            
            result_json = json.loads(result)
            
            if 'output' in result_json and 'choices' in result_json['output']:
                content = result_json['output']['choices'][0]['message']['content']
                
                if isinstance(content, list):
                    text = content[0].get('text', '') if content else ''
                else:
                    text = content
                
                # 尝试解析 JSON，否则直接返回文本
                try:
                    answer_data = self._extract_json(text)
                except:
                    answer_data = {"answer": text}
                
                # 保存对话
                conn = self._get_db()
                cursor = conn.cursor()
                cursor.execute('''
                    INSERT INTO chat_history (session_id, role, content)
                    VALUES (?, ?, ?)
                ''', (session_id, 'user', question))
                cursor.execute('''
                    INSERT INTO chat_history (session_id, role, content)
                    VALUES (?, ?, ?)
                ''', (session_id, 'assistant', json.dumps(answer_data, ensure_ascii=False)))
                conn.commit()
                conn.close()
                
                return {
                    'success': True,
                    'session_id': session_id,
                    'answer': answer_data
                }
            else:
                return {
                    'success': False,
                    'error': 'API 响应异常'
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def get_session(self, session_id: str) -> Optional[Dict]:
        """获取会话详情"""
        conn = self._get_db()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT * FROM document_sessions WHERE session_id = ?
        ''', (session_id,))
        session = cursor.fetchone()
        
        if not session:
            conn.close()
            return None
        
        # 获取对话历史
        cursor.execute('''
            SELECT role, content, created_at FROM chat_history 
            WHERE session_id = ? ORDER BY created_at
        ''', (session_id,))
        history = cursor.fetchall()
        
        conn.close()
        
        return {
            'session_id': session['session_id'],
            'document_type': session['document_type'],
            'document_title': session['document_title'],
            'image_path': session['image_path'],
            'analysis_result': json.loads(session['analysis_result'] or '{}'),
            'mindmap_data': json.loads(session['mindmap_data'] or '{}'),
            'created_at': session['created_at'],
            'chat_history': [
                {'role': h[0], 'content': h[1], 'time': h[2]} for h in history
            ]
        }
    
    def get_sessions(self, limit: int = 20) -> List[Dict]:
        """获取所有会话列表"""
        conn = self._get_db()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT session_id, document_type, document_title, created_at
            FROM document_sessions
            ORDER BY updated_at DESC
            LIMIT ?
        ''', (limit,))
        
        sessions = cursor.fetchall()
        conn.close()
        
        return [
            {
                'session_id': s[0],
                'document_type': s[1],
                'document_title': s[2],
                'created_at': s[3]
            }
            for s in sessions
        ]
    
    def generate_mindmap_from_text(
        self,
        text: str,
        scenario: str = "general",
        session_id: str = None
    ) -> Dict:
        """
        从纯文本生成思维导图（无需图片）
        
        Args:
            text: 用户输入的文本内容
            scenario: 分析场景
            session_id: 会话ID（可选）
            
        Returns:
            分析结果，包含思维导图结构
        """
        logger.info(f"开始从文本生成思维导图: scenario={scenario}, text_length={len(text)}")
        
        # 获取场景配置
        scenario_config = self.SCENARIOS.get(scenario, self.SCENARIOS["general"])
        
        # 构建系统提示词 - 专门针对纯文本输入
        system_prompt = f"""你是一位专业的文档分析和知识结构化专家。请根据用户提供的文本内容，分析并生成思维导图结构。

分析要求：
1. 提取文本的核心主题和关键信息
2. 识别主要概念和它们之间的层级关系
3. 归纳要点，形成清晰的树状结构
4. 确保思维导图覆盖文本的主要内容

输出格式必须是严格的 JSON：
{{
    "document_type": "{scenario}",
    "title": "文档标题",
    "key_content": ["要点1", "要点2", "要点3"],
    "summary": "文档总结",
    "mindmap_structure": {{
        "root": "中心主题",
        "children": [
            {{
                "name": "一级节点",
                "children": [
                    {{"name": "二级节点"}},
                    {{"name": "二级节点"}}
                ]
            }}
        ]
    }}
}}

请确保 mindmap_structure 是一个有效的树形结构，包含 root 和 children 字段。"""
        
        # 构建用户提示
        user_prompt = f"""{system_prompt}

用户输入的文本内容：
---
{text}
---

请分析以上内容，返回 JSON 格式的思维导图结构。"""
        
        # 构建请求 - 纯文本模式，不包含图片
        request_body = json.dumps({
            "model": "qwen-max",  # 使用文本模型，不需要多模态
            "input": {
                "messages": [
                    {
                        "role": "user",
                        "content": user_prompt
                    }
                ]
            },
            "parameters": {
                "result_format": "message"
            }
        })
        
        logger.info(f"请求体大小: {len(request_body.encode('utf-8')) / 1024:.2f} KB")
        
        try:
            conn = http.client.HTTPSConnection("dashscope.aliyuncs.com", timeout=60)
            headers = {
                'Authorization': f'Bearer {self.api_key}',
                'Content-Type': 'application/json'
            }
            
            conn.request(
                "POST",
                "/api/v1/services/aigc/text-generation/generation",
                body=request_body,
                headers=headers
            )
            
            response = conn.getresponse()
            result = response.read().decode('utf-8')
            conn.close()
            
            logger.info(f"API 响应长度: {len(result)} 字符")
            
            # 解析响应
            result_json = json.loads(result)
            
            if 'output' in result_json and 'choices' in result_json['output']:
                content = result_json['output']['choices'][0]['message']['content']
                
                if isinstance(content, list):
                    text_content = content[0].get('text', '') if content else ''
                else:
                    text_content = content
                
                logger.info(f"提取的文本长度: {len(text_content)} 字符")
                
                # 提取 JSON
                analysis = self._extract_json(text_content)
                
                if 'error' in analysis:
                    logger.error(f"JSON 解析失败: {analysis['error']}")
                    return {
                        'success': False,
                        'error': f'JSON 解析失败: {analysis["error"]}',
                        'raw_text': analysis.get('raw_text', '')
                    }
                
                logger.info(f"JSON 解析成功, 包含字段: {list(analysis.keys())}")
                
                # 生成会话ID
                if not session_id:
                    session_id = f"text_session_{datetime.now().strftime('%Y%m%d%H%M%S')}_{hashlib.md5(text.encode()).hexdigest()[:6]}"
                
                # 保存会话（纯文本模式，image_path 为空）
                self._save_text_session(session_id, scenario, analysis, text)
                logger.info(f"会话已保存: session_id={session_id}")
                
                return {
                    'success': True,
                    'session_id': session_id,
                    'analysis': analysis,
                    'mindmap_data': analysis.get('mindmap_structure', {})
                }
            else:
                logger.error(f"API 响应格式异常: {result_json.keys()}")
                return {
                    'success': False,
                    'error': 'API 响应格式异常',
                    'raw_response': result_json
                }
                
        except Exception as e:
            logger.exception(f"生成思维导图失败: {type(e).__name__}: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def _save_text_session(
        self,
        session_id: str,
        document_type: str,
        analysis: Dict,
        source_text: str
    ):
        """保存纯文本会话到数据库"""
        conn = self._get_db()
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT INTO document_sessions 
            (session_id, document_type, document_title, image_path, analysis_result, mindmap_data)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (session_id, document_type,
              analysis.get('title', '未命名文档'),
              '',  # 纯文本模式，无图片路径
              json.dumps(analysis, ensure_ascii=False),
              json.dumps(analysis.get('mindmap_structure', {}), ensure_ascii=False)))
        
        # 保存原始文本作为第一条对话
        cursor.execute('''
            INSERT INTO chat_history (session_id, role, content)
            VALUES (?, ?, ?)
        ''', (session_id, 'user', f"[文本输入] {source_text[:500]}..." if len(source_text) > 500 else f"[文本输入] {source_text}"))
        
        cursor.execute('''
            INSERT INTO chat_history (session_id, role, content)
            VALUES (?, ?, ?)
        ''', (session_id, 'assistant', json.dumps(analysis, ensure_ascii=False)))
        
        conn.commit()
        conn.close()

    def export_to_xmind(self, mindmap_data: Dict, output_path: str = None, title: str = "思维导图") -> str:
        """
        将思维导图导出为 .xmind 文件
        
        Args:
            mindmap_data: 思维导图数据结构
            output_path: 输出文件路径（可选，默认生成临时文件）
            title: 思维导图标题
            
        Returns:
            生成的 .xmind 文件路径
        """
        if not mindmap_data:
            raise ValueError("思维导图数据为空")
        
        # 如果没有指定输出路径，生成临时文件
        if not output_path:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = f"/tmp/mindmap_{timestamp}.xmind"
        
        if XMind_AVAILABLE:
            # 使用 xmind 库生成
            return self._export_with_xmind_lib(mindmap_data, output_path, title)
        else:
            # 使用备用方案生成
            return self._export_with_fallback(mindmap_data, output_path, title)
    
    def _export_with_xmind_lib(self, mindmap_data: Dict, output_path: str, title: str) -> str:
        """使用 xmind 库生成 .xmind 文件"""
        # 创建工作簿
        workbook = xmind.load(output_path)
        sheet = workbook.getPrimarySheet()
        sheet.setTitle(title)
        
        # 获取根主题
        root_topic = sheet.getRootTopic()
        root_topic.setTitle(mindmap_data.get('root', title))
        
        # 递归添加子节点
        def add_children(parent_topic, children_data):
            if not children_data:
                return
            for child in children_data:
                child_topic = parent_topic.addSubTopic()
                child_topic.setTitle(child.get('name', '未命名'))
                # 递归添加子节点
                if 'children' in child and child['children']:
                    add_children(child_topic, child['children'])
        
        # 从根节点开始添加
        children = mindmap_data.get('children', [])
        add_children(root_topic, children)
        
        # 保存文件
        xmind.save(workbook, output_path)
        return output_path
    
    def _export_with_fallback(self, mindmap_data: Dict, output_path: str, title: str) -> str:
        """
        使用备用方案生成 .xmind 文件
        XMind 文件本质上是 ZIP 压缩包，包含 XML 内容
        """
        import xml.etree.ElementTree as ET
        
        # 创建临时目录
        temp_dir = f"/tmp/xmind_temp_{datetime.now().strftime('%Y%m%d%H%M%S')}"
        os.makedirs(temp_dir, exist_ok=True)
        
        # 创建 content.xml
        content_root = ET.Element("xmap-content", {
            "version": "2.0",
            "xmlns": "urn:xmind:xmap:xmlns:content:2.0"
        })
        
        # 创建 sheet
        sheet_elem = ET.SubElement(content_root, "sheet", {"id": "1"})
        title_elem = ET.SubElement(sheet_elem, "title")
        title_elem.text = title
        
        # 创建根主题
        root_topic_elem = ET.SubElement(sheet_elem, "topic", {
            "id": "root",
            "structure-class": "org.xmind.ui.map.unbalanced"
        })
        root_title = ET.SubElement(root_topic_elem, "title")
        root_title.text = mindmap_data.get('root', title)
        
        # 递归添加子节点
        def add_children_xml(parent_elem, children_data, parent_id="root"):
            if not children_data:
                return
            children_elem = ET.SubElement(parent_elem, "children")
            topics_elem = ET.SubElement(children_elem, "topics", {"type": "attached"})
            
            for i, child in enumerate(children_data):
                child_id = f"{parent_id}_{i}"
                child_topic = ET.SubElement(topics_elem, "topic", {"id": child_id})
                child_title = ET.SubElement(child_topic, "title")
                child_title.text = child.get('name', '未命名')
                
                # 递归添加子节点
                if 'children' in child and child['children']:
                    add_children_xml(child_topic, child['children'], child_id)
        
        # 添加根节点的子节点
        children = mindmap_data.get('children', [])
        if children:
            add_children_xml(root_topic_elem, children)
        
        # 写入 content.xml
        content_tree = ET.ElementTree(content_root)
        content_path = os.path.join(temp_dir, "content.xml")
        content_tree.write(content_path, encoding='UTF-8', xml_declaration=True)
        
        # 创建 META-INF/manifest.xml
        meta_inf_dir = os.path.join(temp_dir, "META-INF")
        os.makedirs(meta_inf_dir, exist_ok=True)
        
        manifest_root = ET.Element("manifest", {
            "xmlns": "urn:xmind:xmap:xmlns:manifest:1.0"
        })
        ET.SubElement(manifest_root, "file-entry", {
            "full-path": "content.xml",
            "media-type": "text/xml"
        })
        ET.SubElement(manifest_root, "file-entry", {
            "full-path": "META-INF/",
            "media-type": ""
        })
        ET.SubElement(manifest_root, "file-entry", {
            "full-path": "META-INF/manifest.xml",
            "media-type": "text/xml"
        })
        
        manifest_tree = ET.ElementTree(manifest_root)
        manifest_path = os.path.join(meta_inf_dir, "manifest.xml")
        manifest_tree.write(manifest_path, encoding='UTF-8', xml_declaration=True)
        
        # 打包为 .xmind 文件（ZIP 格式）
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zf.write(file_path, arcname)
        
        # 清理临时目录
        shutil.rmtree(temp_dir)
        
        return output_path

    def generate_mindmap_html(self, mindmap_data: Dict) -> str:
        """生成思维导图的 HTML 代码（使用 Markmap）"""
        if not mindmap_data:
            return "<p>暂无思维导图数据</p>"
        
        # 转换为 Markdown 格式
        def to_markdown(node, level=0):
            indent = "  " * level
            md = f"{indent}- {node.get('name', '未命名')}\n"
            for child in node.get('children', []):
                md += to_markdown(child, level + 1)
            return md
        
        markdown_content = to_markdown(mindmap_data)
        
        html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>思维导图</title>
    <script src="https://cdn.jsdelivr.net/npm/markmap-autoloader@0.17"></script>
    <style>
        body {{ margin: 0; padding: 20px; font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; }}
        #mindmap {{ width: 100%; height: 80vh; }}
        .mindmap-container {{ border: 1px solid #e0e0e0; border-radius: 8px; padding: 20px; }}
    </style>
</head>
<body>
    <div class="mindmap-container">
        <div id="mindmap">
            <pre class="language-markdown">
{markdown_content}
            </pre>
        </div>
    </div>
    <script>
        document.addEventListener('DOMContentLoaded', () => {{
            const el = document.getElementById('mindmap');
            const markdown = el.querySelector('pre').textContent;
            markmap.autoLoader.render(el, markdown);
        }});
    </script>
</body>
</html>'''
        return html

    def generate_ppt_from_mindmap(
        self,
        mindmap_data: Dict,
        chat_history: List[Dict] = None,
        title: str = "演示文稿",
        theme: str = "blue"
    ) -> str:
        """
        将思维导图和对话历史转换为 PPT
        
        Args:
            mindmap_data: 思维导图数据结构
            chat_history: 对话历史记录
            title: PPT 标题
            theme: 主题颜色 (blue, green, purple, orange, dark)
            
        Returns:
            生成的 PPT 文件路径
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 库未安装，无法生成 PPT")
        
        # 主题颜色配置
        themes = {
            "blue": {"primary": RGBColor(0x1E, 0x3A, 0x5F), "accent": RGBColor(0x4A, 0x90, 0xE2)},
            "green": {"primary": RGBColor(0x1B, 0x5E, 0x20), "accent": RGBColor(0x66, 0xBB, 0x6A)},
            "purple": {"primary": RGBColor(0x4A, 0x14, 0x8C), "accent": RGBColor(0x9C, 0x27, 0xB0)},
            "orange": {"primary": RGBColor(0xE6, 0x51, 0x00), "accent": RGBColor(0xFF, 0x98, 0x00)},
            "dark": {"primary": RGBColor(0x21, 0x21, 0x21), "accent": RGBColor(0x61, 0x61, 0x61)},
        }
        colors = themes.get(theme, themes["blue"])
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 1. 添加封面
        self._add_ppt_title_slide(prs, title, "基于 AI 分析生成", colors)
        
        # 2. 添加目录页
        if mindmap_data and mindmap_data.get('children'):
            toc_items = [child.get('name', f'章节 {i+1}') for i, child in enumerate(mindmap_data.get('children', []))]
            self._add_ppt_content_slide(prs, "目录", toc_items, colors)
        
        # 3. 根据思维导图添加内容页
        if mindmap_data:
            root_name = mindmap_data.get('root', '主题')
            children = mindmap_data.get('children', [])
            
            for i, child in enumerate(children):
                section_title = child.get('name', f'章节 {i+1}')
                
                # 添加章节分隔页
                self._add_ppt_section_slide(prs, section_title, colors)
                
                # 收集该章节下的所有要点
                points = []
                sub_children = child.get('children', [])
                
                for sub in sub_children:
                    point = sub.get('name', '')
                    if point:
                        points.append(point)
                        # 添加子要点
                        sub_sub = sub.get('children', [])
                        for ss in sub_sub:
                            sub_point = ss.get('name', '')
                            if sub_point:
                                points.append(f"  • {sub_point}")
                
                # 如果有要点，添加内容页
                if points:
                    # 每页最多显示 6 个要点
                    for j in range(0, len(points), 6):
                        page_points = points[j:j+6]
                        page_title = section_title if j == 0 else f"{section_title} (续)"
                        self._add_ppt_content_slide(prs, page_title, page_points, colors)
                else:
                    # 没有子要点，显示章节标题
                    self._add_ppt_content_slide(prs, section_title, ["本章暂无详细内容"], colors)
        
        # 4. 添加对话分析摘要（如果有）
        if chat_history:
            # 提取关键对话内容
            key_qa = []
            for chat in chat_history[-6:]:  # 最近 6 条对话
                if chat.get('role') == 'user':
                    content = chat.get('content', '')
                    if len(content) > 100:
                        content = content[:100] + '...'
                    key_qa.append(f"Q: {content}")
                elif chat.get('role') == 'assistant':
                    try:
                        content = json.loads(chat.get('content', '{}'))
                        answer = content.get('answer', '')
                        if len(answer) > 100:
                            answer = answer[:100] + '...'
                        if answer:
                            key_qa.append(f"A: {answer}")
                    except:
                        pass
            
            if key_qa:
                self._add_ppt_section_slide(prs, "对话分析摘要", colors)
                self._add_ppt_content_slide(prs, "关键问答", key_qa[:6], colors)
        
        # 5. 添加结束页
        self._add_ppt_end_slide(prs, "谢谢观看", colors)
        
        # 保存文件
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_path = f"/tmp/ppt_{timestamp}.pptx"
        prs.save(output_path)
        
        return output_path
    
    def _add_ppt_title_slide(self, prs, title, subtitle, colors):
        """添加 PPT 封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors["primary"]
        bg.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_ppt_content_slide(self, prs, title, points, colors):
        """添加 PPT 内容页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors["accent"]
        bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 处理缩进
            if point.startswith("  "):
                p.text = point.strip()
                p.level = 1
                p.font.size = Pt(16)
            else:
                p.text = f"• {point}"
                p.font.size = Pt(20)
            
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(12)
    
    def _add_ppt_section_slide(self, prs, title, colors):
        """添加 PPT 章节分隔页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors["accent"]
        bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_ppt_end_slide(self, prs, message, colors):
        """添加 PPT 结束页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors["primary"]
        bg.line.fill.background()
        
        msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = message
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


# 便捷函数
def analyze_document(image_path: str, question: str = None, scenario: str = "general") -> Dict:
    """便捷函数：分析文档"""
    generator = DocumentMindMapGenerator()
    return generator.analyze_document(image_path, question, scenario)


def chat_with_document(session_id: str, question: str) -> Dict:
    """便捷函数：与文档对话"""
    generator = DocumentMindMapGenerator()
    return generator.chat(session_id, question)


def generate_mindmap_from_text(text: str, scenario: str = "general") -> Dict:
    """便捷函数：从纯文本生成思维导图"""
    generator = DocumentMindMapGenerator()
    return generator.generate_mindmap_from_text(text, scenario)


if __name__ == '__main__':
    import argparse
    
    parser = argparse.ArgumentParser(description='文档分析与思维导图生成器')
    parser.add_argument('image', help='文档图片路径')
    parser.add_argument('--question', '-q', help='分析问题')
    parser.add_argument('--scenario', '-s', default='general', 
                       choices=['financial_report', 'flowchart', 'architecture', 'mindmap', 'general'],
                       help='分析场景')
    parser.add_argument('--api-key', help='API Key')
    
    args = parser.parse_args()
    
    generator = DocumentMindMapGenerator({'api_key': args.api_key})
    result = generator.analyze_document(args.image, args.question, args.scenario)
    
    if result['success']:
        print(f"\n✅ 分析成功！")
        print(f"会话ID: {result['session_id']}")
        print(f"\n分析结果:")
        print(json.dumps(result['analysis'], ensure_ascii=False, indent=2))
    else:
        print(f"\n❌ 分析失败: {result.get('error')}")